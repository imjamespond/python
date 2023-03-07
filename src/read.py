from pptx import Presentation

""" 
slide 幻灯片页
shape 形状
paragraph 段落
run 文字块 
"""

ppt = Presentation("test.pptx")

# slides 得到一个列表,包含每个列表slide
for slide in ppt.slides:
  print(slide)
  
  # slide.shapes 形状
  for shape in slide.shapes:
    print(shape)

    # 判断shape内的文字，并进行输出
    if shape.has_text_frame:
      text_frame = shape.text_frame
      print(text_frame.text)

      # 寻找paragra
      for paragraph in text_frame.paragraphs:
        print(paragraph.text)